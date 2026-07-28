from __future__ import annotations
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from core.metrics import kpis, monthly_summary, demand_summary, stop_frequency, punctuality_summary, time_statistics
from ai.conclusions import generate_conclusions

BLUE = RGBColor(18, 61, 118)
LIGHT = RGBColor(238, 243, 249)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(31, 41, 55)


def _title(slide, title, subtitle=""):
    shape = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(0.9))
    shape.fill.solid(); shape.fill.fore_color.rgb = BLUE; shape.line.fill.background()
    box = slide.shapes.add_textbox(Inches(.45), Inches(.18), Inches(12.3), Inches(.55))
    p = box.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(25); p.font.bold = True; p.font.color.rgb = WHITE
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(.5), Inches(1.02), Inches(12.1), Inches(.4))
        p = sub.text_frame.paragraphs[0]; p.text = subtitle; p.font.size = Pt(12); p.font.color.rgb = DARK


def _metric(slide, x, y, w, h, label, value):
    sh = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = LIGHT; sh.line.color.rgb = BLUE
    tf = sh.text_frame; tf.clear()
    p = tf.paragraphs[0]; p.text = str(value); p.alignment = PP_ALIGN.CENTER; p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = BLUE
    p2 = tf.add_paragraph(); p2.text = label; p2.alignment = PP_ALIGN.CENTER; p2.font.size = Pt(10); p2.font.color.rgb = DARK


def _table(slide, df, x=.5, y=1.5, w=12.3, h=5.3, max_rows=12):
    view = df.head(max_rows).copy()
    if view.empty:
        return
    rows, cols = len(view) + 1, len(view.columns)
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for j, col in enumerate(view.columns):
        cell = table.cell(0, j); cell.text = str(col); cell.fill.solid(); cell.fill.fore_color.rgb = BLUE
        for p in cell.text_frame.paragraphs: p.font.bold = True; p.font.color.rgb = WHITE; p.font.size = Pt(9)
    for i, (_, row) in enumerate(view.iterrows(), start=1):
        for j, value in enumerate(row):
            text = f"{value:.1%}" if isinstance(value, float) and 0 <= value <= 1 else (f"{value:.1f}" if isinstance(value, float) else str(value))
            table.cell(i, j).text = text
            for p in table.cell(i, j).text_frame.paragraphs: p.font.size = Pt(8); p.font.color.rgb = DARK


def _bar_chart(slide, title, categories, series, x=.6, y=1.5, w=12.0, h=5.2):
    data = ChartData(); data.categories = [str(x) for x in categories]
    for name, values in series.items(): data.add_series(name, [0 if v is None else float(v) for v in values])
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), data).chart
    chart.has_title = True; chart.chart_title.text_frame.text = title; chart.has_legend = len(series) > 1


def build_ppt(df) -> bytes:
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    m = kpis(df)

    slide = prs.slides.add_slide(blank); _title(slide, "KOA Analytics — Informe ejecutivo", "Estudios de tiempos, puntualidad y optimización de paraderos")
    metrics = [
        ("Registros", m["registros"]), ("Usuarios", m["usuarios"]), ("Efectivos", m["efectivos"]),
        ("Puntualidad general", f"{m['puntualidad_general']:.1%}"), ("Puntualidad mañana", f"{m['puntualidad_manana']:.1%}" if m['puntualidad_manana']==m['puntualidad_manana'] else "N/D"),
        ("Puntualidad tarde", f"{m['puntualidad_tarde']:.1%}" if m['puntualidad_tarde']==m['puntualidad_tarde'] else "N/D"),
        ("Tiempo promedio", f"{m['tiempo_efectivo_promedio']:.1f} min"), ("P90", f"{m['tiempo_efectivo_p90']:.1f} min"),
    ]
    for idx, (label, value) in enumerate(metrics): _metric(slide, .55+(idx%4)*3.15, 1.55+(idx//4)*2.05, 2.8, 1.55, label, value)

    monthly = monthly_summary(df)
    if not monthly.empty:
        slide = prs.slides.add_slide(blank); _title(slide, "Evolución mensual")
        _bar_chart(slide, "Usuarios por mes", monthly["mes"], {"Usuarios": monthly["usuarios"]})

    demand = demand_summary(df)
    if not demand.empty:
        slide = prs.slides.add_slide(blank); _title(slide, "Demanda por recorrido")
        _bar_chart(slide, "Usuarios y recorridos efectivos", demand["recorrido"], {"Usuarios": demand["usuarios"], "Efectivos": demand["efectivos"]})

    stops = stop_frequency(df)
    if not stops.empty:
        slide = prs.slides.add_slide(blank); _title(slide, "Utilización de paraderos")
        _table(slide, stops[["paradero", "usos_efectivos", "frecuencia_efectiva", "usuarios_asociados", "tiempo_promedio"]])

    punctual = punctuality_summary(df)
    if not punctual.empty:
        slide = prs.slides.add_slide(blank); _title(slide, "Puntualidad por jornada y recorrido")
        _table(slide, punctual[["jornada", "recorrido", "registros_validos", "puntualidad_oficial", "anticipadas", "retraso_mayor_5"]])

    stats = time_statistics(df)
    if not stats.empty:
        slide = prs.slides.add_slide(blank); _title(slide, "Estudio estadístico de tiempos")
        _table(slide, stats)

    slide = prs.slides.add_slide(blank); _title(slide, "Hallazgos y recomendaciones")
    box = slide.shapes.add_textbox(Inches(.7), Inches(1.45), Inches(12), Inches(5.3)); tf = box.text_frame; tf.word_wrap = True
    conclusions = generate_conclusions(df) or ["No existe evidencia suficiente con los filtros seleccionados."]
    for idx, item in enumerate(conclusions):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph(); p.text = f"• {item}"; p.font.size = Pt(16); p.space_after = Pt(10); p.font.color.rgb = DARK
    p = tf.add_paragraph(); p.text = "• Las propuestas de mapa y simulación deben validarse mediante visita de campo y prueba piloto."; p.font.size = Pt(16); p.font.color.rgb = DARK

    output = BytesIO(); prs.save(output); return output.getvalue()
